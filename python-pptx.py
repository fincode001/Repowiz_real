# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# --- 프레젠테이션 객체 생성 ---
prs = Presentation()
# 슬라이드 레이아웃 정의 (기본 템플릿 사용)
title_slide_layout = prs.slide_layouts[0] # 제목 슬라이드
bullet_slide_layout = prs.slide_layouts[1] # 제목 및 내용 (불렛)
title_only_layout = prs.slide_layouts[5] # 제목만

# --- 슬라이드 생성 ---

# 슬라이드 1: 제목 슬라이드
slide1 = prs.slides.add_slide(title_slide_layout)
title1 = slide1.shapes.title
subtitle1 = slide1.placeholders[1]
title1.text = "성공적인 신규 사업 추진을 위한 핵심: 사업타당성 분석 프레임워크"
subtitle1.text = "보고 대상: 최고 경영진\n작성일: 2025. 4. 16." # 날짜 동적 삽입

# 슬라이드 2: Executive Summary
slide2 = prs.slides.add_slide(bullet_slide_layout)
title2 = slide2.shapes.title
body2 = slide2.shapes.placeholders[1]
title2.text = "Executive Summary"
tf2 = body2.text_frame
tf2.text = "사업타당성 분석의 중요성 및 핵심 검토 영역 요약:"
p2_1 = tf2.add_paragraph()
p2_1.text = "신규 사업 성공률 제고를 위해 객관적 사업타당성 분석은 필수적입니다."
p2_1.level = 1
p2_2 = tf2.add_paragraph()
p2_2.text = "주요 검토 영역: 시장 매력도, 기술 구현 가능성, 재무적 수익성, 내부 수행 역량 등"
p2_2.level = 1
# 실제 요약 내용 추가 (시뮬레이션)
p2_3 = tf2.add_paragraph()
p2_3.text = "(샘플)신규사업타당성 검토.txt 요약...\n사업타당성 분석은 성공 확률을 높이는 핵심... 객관적 분석과 시장성, 기술성, 재무성 검토 필요......" # 줄바꿈 처리
p2_3.level = 1

# 슬라이드 3: 사업타당성 분석의 필요성
slide3 = prs.slides.add_slide(bullet_slide_layout)
title3 = slide3.shapes.title
body3 = slide3.shapes.placeholders[1]
title3.text = "왜 사업타당성 분석이 중요한가?"
tf3 = body3.text_frame
tf3.add_paragraph().text = "실패 위험 감소: 주관적 판단 배제, 객관적 데이터 기반 의사결정"
tf3.add_paragraph().text = "성공 확률 증대: 시장 기회 포착 및 잠재 리스크 사전 식별"
tf3.add_paragraph().text = "자원 배분 효율화: 제한된 자원을 성공 가능성 높은 사업에 집중"
tf3.add_paragraph().text = "투자 유치 및 이해관계자 설득력 강화"

# 슬라이드 4: 핵심 분석 영역 (시장성)
slide4 = prs.slides.add_slide(bullet_slide_layout)
title4 = slide4.shapes.title
body4 = slide4.shapes.placeholders[1]
title4.text = "핵심 분석 영역 1: 시장성 (Marketability)"
tf4 = body4.text_frame
tf4.add_paragraph().text = "시장 규모 및 성장성: 진입하려는 시장은 충분히 크고 성장하고 있는가?"
tf4.add_paragraph().text = "타겟 고객 정의: 우리의 핵심 고객은 누구이며, 그들의 니즈는 무엇인가?"
tf4.add_paragraph().text = "경쟁 환경 분석: 주요 경쟁자는 누구이며, 우리의 경쟁 우위는 무엇인가?"
tf4.add_paragraph().text = "가격 및 유통 전략: 적절한 가격 정책과 효과적인 유통 채널 확보 방안은?"

# 슬라이드 5: 핵심 분석 영역 (기술성)
slide5 = prs.slides.add_slide(bullet_slide_layout)
title5 = slide5.shapes.title
body5 = slide5.shapes.placeholders[1]
title5.text = "핵심 분석 영역 2: 기술성 (Technical Feasibility)"
tf5 = body5.text_frame
tf5.add_paragraph().text = "제품/서비스 구현 가능성: 보유 기술 또는 확보 가능한 기술로 구현 가능한가?"
tf5.add_paragraph().text = "생산/운영 안정성: 안정적인 생산 및 서비스 운영 체계 구축이 가능한가?"
tf5.add_paragraph().text = "핵심 기술 확보 및 보호: 핵심 기술은 확보되었는가? 지식재산권 확보 방안은?"
tf5.add_paragraph().text = "기술 변화 대응력: 미래 기술 변화에 유연하게 대응할 수 있는가?"

# 슬라이드 6: 핵심 분석 영역 (재무성)
slide6 = prs.slides.add_slide(bullet_slide_layout)
title6 = slide6.shapes.title
body6 = slide6.shapes.placeholders[1]
title6.text = "핵심 분석 영역 3: 재무성 (Financial Viability)"
tf6 = body6.text_frame
tf6.add_paragraph().text = "수익성 예측: 예상 매출과 비용 구조는 어떠하며, 충분한 수익 확보가 가능한가?"
tf6.add_paragraph().text = "손익분기점(BEP) 분석: 언제 BEP 달성이 가능하며, 그 규모는 어느 정도인가?"
tf6.add_paragraph().text = "투자수익률(ROI) 분석: 투자 대비 기대 수익은 매력적인가?"
tf6.add_paragraph().text = "자금 조달 계획: 필요한 자금 규모는 얼마이며, 안정적인 조달 방안은 마련되었는가?"

# 슬라이드 7: 분석 프로세스 제안
slide7 = prs.slides.add_slide(title_only_layout) # 내용 대신 시각화 제안
title7 = slide7.shapes.title
title7.text = "합리적인 분석 프로세스"
# 시각화 제안
left, top, width, height = Inches(1), Inches(1.5), Inches(8), Inches(4)
textbox = slide7.shapes.add_textbox(left, top, width, height)
tf7 = textbox.text_frame
p7 = tf7.add_paragraph()
p7.text = "[시각화 제안: 아이템 탐색 → 예비 분석 → 상세 분석(시장/기술/재무) → 종합 평가 → 의사결정 과정을 보여주는 플로우차트 삽입]"
p7.font.size = Pt(18)
p7.alignment = PP_ALIGN.CENTER

# 슬라이드 8: 성공적인 분석을 위한 제언
slide8 = prs.slides.add_slide(bullet_slide_layout)
title8 = slide8.shapes.title
body8 = slide8.shapes.placeholders[1]
title8.text = "성공적인 분석을 위한 제언"
tf8 = body8.text_frame
tf8.add_paragraph().text = "데이터 기반 객관성 확보: 직감이나 편견 대신 데이터를 기반으로 판단"
tf8.add_paragraph().text = "최신 트렌드 반영: ESG, 디지털 전환 등 변화하는 외부 환경 요인 고려 (디지털 전환, ESG 중요성 부각...)"
tf8.add_paragraph().text = "위험 관리 및 시나리오 플래닝: 잠재적 위험 요소를 식별하고 대응 계획 수립 (창업자 역량 평가 방법, 정보 신뢰성...)"
tf8.add_paragraph().text = "내/외부 전문가 적극 활용: 필요시 외부 컨설팅, 내부 관련 부서 협업"

# 슬라이드 9: 결론 및 다음 단계
slide9 = prs.slides.add_slide(bullet_slide_layout)
title9 = slide9.shapes.title
body9 = slide9.shapes.placeholders[1]
title9.text = "결론 및 다음 단계"
tf9 = body9.text_frame
tf9.add_paragraph().text = "체계적인 사업타당성 분석은 신규 사업 성공의 초석입니다."
tf9.add_paragraph().text = "본 프레임워크를 활용하여 향후 신규 사업 검토의 객관성과 효율성을 제고할 것을 제안합니다."
tf9.add_paragraph().text = "다음 단계: 전사적 사업타당성 검토 프로세스 표준화 논의"


# --- PPTX 파일 저장 ---
file_path = 'ReportWiz_Generated_사업타당성검토_임원보고.pptx'
prs.save(file_path)

print(f"PPTX 파일이 성공적으로 생성되었습니다: {file_path}")